import base64
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PptxParser:
    async def parse(self, filepath: str) -> dict:
        prs = Presentation(filepath)
        slides_text: list[str] = []
        images: list[dict] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text)

                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        slide_parts.append(" | ".join(cells))

                # Extract images from shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        raw = shape.image.blob
                        content_type = shape.image.content_type or "image/png"
                        b64 = base64.b64encode(raw).decode("ascii")
                        images.append({"base64": b64, "mediaType": content_type})
                    except Exception:
                        pass

            slides_text.append("\n".join(slide_parts))

        text = "\n\n".join(slides_text)
        result: dict = {"text": text}
        if images:
            result["images"] = images
            result["extractMethod"] = "hybrid" if text.strip() else "vision"
        else:
            result["extractMethod"] = "text"

        return result
